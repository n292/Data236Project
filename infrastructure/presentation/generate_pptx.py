"""
LinkedIn Platform - Group Presentation PPTX Generator
Covers all 6 required sections from the project spec.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches, Pt
import copy

# ── Brand palette ──────────────────────────────────────────────────────────────
LI_BLUE   = RGBColor(0x0A, 0x66, 0xC2)   # LinkedIn primary blue
LI_DARK   = RGBColor(0x00, 0x1B, 0x48)   # Dark navy
LI_WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
LI_LIGHT  = RGBColor(0xF3, 0xF2, 0xEF)   # LinkedIn background gray
LI_GREEN  = RGBColor(0x05, 0x7A, 0x55)
LI_ORANGE = RGBColor(0xE8, 0x6A, 0x23)
LI_GRAY   = RGBColor(0x66, 0x66, 0x66)
LI_ACCENT = RGBColor(0x70, 0xB5, 0xF9)   # Light blue accent

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

BLANK = prs.slide_layouts[6]   # completely blank

# ─── Helper functions ──────────────────────────────────────────────────────────

def add_rect(slide, l, t, w, h, fill=LI_BLUE, alpha=None):
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.line.fill.background()
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    return shape

def add_text(slide, text, l, t, w, h, size=18, bold=False, color=LI_WHITE,
             align=PP_ALIGN.LEFT, wrap=True, italic=False):
    txBox = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    txBox.word_wrap = wrap
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox

def add_bullet_box(slide, items, l, t, w, h, size=14, title=None, title_color=LI_BLUE,
                   bullet_color=LI_DARK, bg=None, pad=0.15):
    if bg:
        add_rect(slide, l, t, w, h, fill=bg)
    if title:
        add_text(slide, title, l+pad, t+0.05, w-pad*2, 0.4, size=size+2,
                 bold=True, color=title_color)
        t += 0.42
        h -= 0.42
    txBox = slide.shapes.add_textbox(Inches(l+pad), Inches(t+0.05), Inches(w-pad*2), Inches(h-0.1))
    txBox.word_wrap = True
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(3)
        run = p.add_run()
        run.text = item
        run.font.size = Pt(size)
        run.font.color.rgb = bullet_color
    return txBox

def slide_header(slide, title, subtitle=None):
    """Blue top bar with title."""
    add_rect(slide, 0, 0, 13.33, 1.15, fill=LI_DARK)
    add_rect(slide, 0, 1.15, 13.33, 0.06, fill=LI_BLUE)
    add_text(slide, title, 0.4, 0.12, 12, 0.7, size=28, bold=True, color=LI_WHITE)
    if subtitle:
        add_text(slide, subtitle, 0.4, 0.75, 12, 0.38, size=14, color=LI_ACCENT)

def slide_footer(slide):
    add_rect(slide, 0, 7.2, 13.33, 0.3, fill=LI_DARK)
    add_text(slide, "DATA 236 · LinkedIn Platform  |  Group 6", 0.3, 7.22, 8, 0.26,
             size=9, color=LI_ACCENT)
    add_text(slide, "2024–25", 12.0, 7.22, 1.2, 0.26, size=9,
             color=LI_ACCENT, align=PP_ALIGN.RIGHT)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 – Title / Cover
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, 13.33, 7.5, fill=LI_DARK)
add_rect(slide, 0, 0, 13.33, 0.08, fill=LI_BLUE)
add_rect(slide, 0, 7.42, 13.33, 0.08, fill=LI_BLUE)

# Decorative circles
for cx, cy, sz, alpha in [(11.5,1.2,2.2,LI_BLUE),(12.8,0.3,1.6,LI_ACCENT),(0.5,6.5,1.2,LI_BLUE)]:
    c = slide.shapes.add_shape(9, Inches(cx), Inches(cy), Inches(sz), Inches(sz))
    c.fill.solid(); c.fill.fore_color.rgb = alpha
    c.line.fill.background()

add_text(slide, "🔗", 0.5, 1.6, 1.2, 1.0, size=52, color=LI_BLUE)
add_text(slide, "LinkedIn Platform", 1.6, 1.55, 10, 1.0, size=44, bold=True, color=LI_WHITE)
add_text(slide, "Scalable Microservices Architecture with AI Orchestration",
         1.6, 2.55, 10, 0.55, size=20, color=LI_ACCENT)

add_rect(slide, 1.6, 3.3, 9.5, 0.04, fill=LI_BLUE)

add_text(slide, "Group 6", 1.6, 3.5, 5, 0.45, size=20, bold=True, color=LI_WHITE)
add_text(slide, "Dipin Jassal  ·  Sarvesh Reshimwale  ·  Sammruddhi  ·  Anushka Khadatkar  ·  Rajesh  ·  Bhavya  ·  Shashira  ·  Nikhil",
         1.6, 3.95, 10, 0.4, size=13, color=LI_ACCENT)
add_text(slide, "DATA 236 – Advanced Database Systems  ·  Spring 2025",
         1.6, 4.4, 10, 0.4, size=14, color=LI_GRAY)

tags = [("MySQL", 1.6), ("MongoDB", 3.1), ("Redis", 4.5), ("Kafka", 5.8), ("FastAPI", 7.1), ("React", 8.5)]
for tag, x in tags:
    add_rect(slide, x, 5.2, 1.25, 0.38, fill=LI_BLUE)
    add_text(slide, tag, x+0.05, 5.22, 1.15, 0.34, size=12, bold=True,
             color=LI_WHITE, align=PP_ALIGN.CENTER)

add_text(slide, "78,985 jobs  ·  10,000+ members  ·  8 microservices  ·  14 Kafka topics",
         1.6, 5.8, 10, 0.4, size=13, color=LI_ACCENT, italic=True)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 – Group Details
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, 13.33, 7.5, fill=LI_LIGHT)
slide_header(slide, "1 · Group & Team Details")
slide_footer(slide)

members = [
    "Dipin Jassal",
    "Sarvesh Reshimwale",
    "Sammruddhi",
    "Anushka Khadatkar",
    "Rajesh",
    "Bhavya",
    "Shashira",
    "Nikhil",
]

for i, name in enumerate(members):
    col = i % 4
    row = i // 4
    x = 0.4 + col * 3.15
    y = 1.45 + row * 2.3
    add_rect(slide, x, y, 2.95, 2.0, fill=LI_WHITE)
    add_rect(slide, x, y, 0.08, 2.0, fill=LI_BLUE)
    # Avatar circle
    av = slide.shapes.add_shape(9, Inches(x+0.28), Inches(y+0.32), Inches(0.95), Inches(0.95))
    av.fill.solid(); av.fill.fore_color.rgb = LI_BLUE; av.line.fill.background()
    initials = "".join(w[0] for w in name.split()[:2]).upper()
    add_text(slide, initials, x+0.28, y+0.38, 0.95, 0.6, size=20, bold=True,
             color=LI_WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, name, x+0.2, y+1.38, 2.6, 0.52, size=13, bold=True,
             color=LI_DARK, align=PP_ALIGN.CENTER)

add_text(slide, "PROJECT:  LinkedIn-Like Professional Network  ·  Microservices + AI Orchestration",
         0.5, 6.75, 12, 0.38, size=12, color=LI_GRAY, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 – Database Schema
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, 13.33, 7.5, fill=LI_LIGHT)
slide_header(slide, "2 · Database Schema", "MySQL (transactional) + MongoDB (document) + Redis (cache)")
slide_footer(slide)

# MySQL tables - left column
mysql_tables = [
    ("members",      ["member_id (PK)", "first_name, last_name", "email (UNIQUE)", "headline, about_summary",
                      "experience_json, education_json", "skills_json", "connections_count", "created_at, updated_at"]),
    ("job_postings", ["job_id (PK, CHAR 36)", "company_id, recruiter_id", "title, description",
                      "employment_type, location", "remote ENUM(onsite/remote/hybrid)", "skills_required JSON",
                      "salary_min, salary_max", "status ENUM(open/closed)", "views_count, applicants_count",
                      "FULLTEXT idx (title, description)"]),
    ("applications", ["application_id (PK, UUID)", "job_id → job_postings", "member_id → members",
                      "resume_file_name, resume_file_path", "cover_letter", "status ENUM(submitted→accepted/rejected)",
                      "recruiter_note", "UNIQUE(job_id, member_id)"]),
]
mysql_colors = [LI_BLUE, RGBColor(0x05, 0x5A, 0xA0), RGBColor(0x02, 0x4A, 0x90)]

for i, (tname, cols) in enumerate(mysql_tables):
    y = 1.3 + i * 1.95
    add_rect(slide, 0.3, y, 4.1, 1.82, fill=mysql_colors[i])
    add_text(slide, f"▪ {tname}", 0.45, y+0.08, 3.8, 0.35, size=13, bold=True, color=LI_WHITE)
    for j, col in enumerate(cols[:5]):
        add_text(slide, f"  {col}", 0.45, y+0.42+j*0.24, 3.8, 0.24, size=9.5, color=LI_WHITE)

# More MySQL tables - middle column
extra_mysql = [
    ("saved_jobs",        ["user_id (PK)", "job_id (PK)", "saved_at DATETIME", "INDEX(user_id), INDEX(job_id)"]),
    ("connections",       ["connection_id (PK)", "requester_id, receiver_id", "status ENUM(pending/accepted/rejected)",
                           "UNIQUE(requester_id, receiver_id)", "INDEX(requester), INDEX(receiver)"]),
    ("processed_events",  ["idempotency_key (PK)", "event_type, trace_id", "entity_type, entity_id",
                           "processed_at", "Kafka at-least-once dedup"]),
]
for i, (tname, cols) in enumerate(extra_mysql):
    y = 1.3 + i * 1.95
    add_rect(slide, 4.65, y, 3.6, 1.82, fill=RGBColor(0x1A, 0x4A, 0x7C))
    add_text(slide, f"▪ {tname}", 4.78, y+0.08, 3.3, 0.35, size=12, bold=True, color=LI_WHITE)
    for j, col in enumerate(cols[:5]):
        add_text(slide, f"  {col}", 4.78, y+0.42+j*0.24, 3.3, 0.24, size=9.5, color=LI_WHITE)

# MongoDB / Redis - right column
add_rect(slide, 8.5, 1.3, 4.55, 1.82, fill=LI_GREEN)
add_text(slide, "▪ MongoDB: analytics.events", 8.62, 1.38, 4.3, 0.35, size=12, bold=True, color=LI_WHITE)
for j, col in enumerate(["idempotency_key (dedup)", "event_type, trace_id, actor_id",
                          "entity_type, entity_id", "payload (flexible JSON)", "14 Kafka topics ingested", "_topic, _ingested_at"]):
    add_text(slide, f"  {col}", 8.62, 1.72+j*0.24, 4.3, 0.24, size=9.5, color=LI_WHITE)

add_rect(slide, 8.5, 3.25, 4.55, 1.82, fill=RGBColor(0x03, 0x5C, 0x40))
add_text(slide, "▪ MongoDB: messaging", 8.62, 3.33, 4.3, 0.35, size=12, bold=True, color=LI_WHITE)
for j, col in enumerate(["threads: thread_id, participants[]", "  last_message_at, message_count",
                          "messages: message_id, thread_id", "  sender_id, message_text", "  status ENUM(sent/delivered/read)", "  idempotency_key (dedup)"]):
    add_text(slide, f"  {col}", 8.62, 3.67+j*0.24, 4.3, 0.24, size=9.5, color=LI_WHITE)

add_rect(slide, 8.5, 5.2, 4.55, 1.82, fill=LI_ORANGE)
add_text(slide, "▪ Redis Cache", 8.62, 5.28, 4.3, 0.35, size=12, bold=True, color=LI_WHITE)
for j, col in enumerate(["search results → TTL 60 s", "job detail → TTL 10 s",
                          "Key: search:{sha256(params)}", "Key: job:{job_id}", "Invalidated on: job update/close", "Reduces MySQL load by ~75%"]):
    add_text(slide, f"  {col}", 8.62, 5.62+j*0.24, 4.3, 0.24, size=9.5, color=LI_WHITE)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 – System Architecture Design Diagram
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, 13.33, 7.5, fill=RGBColor(0xF8, 0xF9, 0xFA))
slide_header(slide, "3 · System Architecture Design Diagram", "8 microservices · 3 databases · Redis cache · Kafka event bus")
slide_footer(slide)

# Draw architecture layers
# Layer labels
for lx, ly, lw, lh, ltext, lcolor in [
    (0.15, 1.25, 1.0, 5.8, "CLIENT\nLAYER", LI_BLUE),
    (1.35, 1.25, 1.0, 5.8, "API\nGATEWAY", LI_DARK),
    (2.55, 1.25, 5.8, 5.8, "MICROSERVICES", LI_BLUE),
    (8.55, 1.25, 2.35, 5.8, "MESSAGING\nLAYER", LI_GREEN),
    (11.1, 1.25, 2.1, 5.8, "DATA\nSTORES", LI_ORANGE),
]:
    add_rect(slide, lx, ly, lw, 0.28, fill=lcolor)
    add_text(slide, ltext, lx+0.05, ly+0.0, lw-0.1, 0.28, size=8, bold=True,
             color=LI_WHITE, align=PP_ALIGN.CENTER)

# Client box
add_rect(slide, 0.15, 1.55, 1.0, 1.2, fill=LI_BLUE)
add_text(slide, "React\nFrontend\n:5173", 0.2, 1.6, 0.9, 1.1, size=9, color=LI_WHITE, align=PP_ALIGN.CENTER)

add_rect(slide, 0.15, 2.95, 1.0, 0.7, fill=RGBColor(0x40, 0x80, 0xC0))
add_text(slide, "Recruiter\nPortal", 0.2, 2.97, 0.9, 0.65, size=9, color=LI_WHITE, align=PP_ALIGN.CENTER)

# Vite Proxy / API Gateway
add_rect(slide, 1.35, 1.55, 1.05, 2.1, fill=LI_DARK)
add_text(slide, "Vite\nProxy\n/api/*", 1.38, 1.58, 0.95, 2.0, size=9, color=LI_WHITE, align=PP_ALIGN.CENTER)

# Arrow from client to gateway
ar = slide.shapes.add_connector(1, Inches(1.15), Inches(2.0), Inches(1.35), Inches(2.0))
ar.line.color.rgb = LI_BLUE; ar.line.width = Pt(2)

# Microservices boxes
svcs = [
    ("Profile Svc\n:8002\n(FastAPI)", 2.55, 1.55, 1.65, LI_BLUE),
    ("Job Svc\n:3002\n(Node.js)", 4.4,  1.55, 1.65, LI_DARK),
    ("Application\nSvc :5003\n(Node.js)", 6.25, 1.55, 1.65, LI_BLUE),
    ("Connection\nSvc :3005\n(Node.js)", 2.55, 3.1,  1.65, RGBColor(0x03, 0x5C, 0x90)),
    ("Messaging\nSvc :3004\n(Node.js)", 4.4,  3.1,  1.65, RGBColor(0x03, 0x5C, 0x90)),
    ("Analytics\nSvc :4000\n(Node.js)", 6.25, 3.1,  1.65, RGBColor(0x03, 0x5C, 0x90)),
    ("AI Service\n:8005\n(FastAPI)", 2.55, 4.65, 1.65, LI_GREEN),
    ("Redis\nCache\n:6379", 4.4,  4.65, 1.65, LI_ORANGE),
]
for label, sx, sy, sw, sc in svcs:
    add_rect(slide, sx, sy, sw, 1.3, fill=sc)
    add_text(slide, label, sx+0.07, sy+0.05, sw-0.14, 1.2, size=9, bold=True,
             color=LI_WHITE, align=PP_ALIGN.CENTER)

# Kafka bus
add_rect(slide, 8.55, 1.55, 2.25, 5.3, fill=RGBColor(0x12, 0x27, 0x40))
add_text(slide, "KAFKA\nEvent Bus\n:9092", 8.6, 1.6, 2.15, 0.8, size=10, bold=True,
         color=LI_ACCENT, align=PP_ALIGN.CENTER)
topics = ["job.created", "job.closed", "job.viewed", "job.saved",
          "application.submitted", "application.status_updated",
          "connection.requested", "connection.accepted",
          "message.sent", "member.created", "member.updated",
          "profile.viewed", "ai.requests", "ai.results"]
for j, t in enumerate(topics):
    add_rect(slide, 8.65, 2.5+j*0.32, 2.1, 0.28, fill=RGBColor(0x1A, 0x40, 0x70))
    add_text(slide, t, 8.68, 2.51+j*0.32, 2.05, 0.26, size=7.5, color=LI_ACCENT)

# Data stores
ds = [
    ("MySQL\ndata236\nlinkedin_sim\napp_db\nconnections", 11.1, 1.55, LI_BLUE),
    ("MongoDB\nanalytics\nmessaging\nai_db", 11.1, 3.55, LI_GREEN),
    ("Redis\nSearch cache\nJob cache\nTTL 10-60s", 11.1, 5.35, LI_ORANGE),
]
for label, dx, dy, dc in ds:
    add_rect(slide, dx, dy, 2.05, 1.65, fill=dc)
    add_text(slide, label, dx+0.08, dy+0.08, 1.9, 1.5, size=9, color=LI_WHITE, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 – Agent Architecture Design Diagram
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, 13.33, 7.5, fill=RGBColor(0x0D, 0x1B, 0x2A))
slide_header(slide, "4 · Agent Architecture & Kafka Topic Design", "AI Orchestration · Hiring Assistant · Career Coach")
slide_footer(slide)

# AI pipeline flow
steps = [
    ("Recruiter\nSubmits Task", 0.3, 2.2, 1.6),
    ("POST\n/ai/submit-task", 2.1, 2.2, 1.6),
    ("Kafka\nai.requests", 3.9, 2.2, 1.6),
    ("AI Service\nWorker", 5.7, 2.2, 1.6),
    ("MongoDB\nStore Task", 7.5, 2.2, 1.6),
    ("Kafka\nai.results", 9.3, 2.2, 1.6),
    ("Analytics\nConsumer", 11.1, 2.2, 1.6),
]
colors = [RGBColor(0x0A,0x66,0xC2), RGBColor(0x05,0x5A,0xA0), RGBColor(0x8B,0x00,0x00),
          LI_GREEN, RGBColor(0x03,0x5C,0x40), RGBColor(0x8B,0x45,0x00), RGBColor(0x5C,0x03,0x5C)]
for i, (label, sx, sy, sw) in enumerate(steps):
    add_rect(slide, sx, sy, sw, 1.1, fill=colors[i])
    add_text(slide, label, sx+0.05, sy+0.1, sw-0.1, 0.9, size=10, bold=True,
             color=LI_WHITE, align=PP_ALIGN.CENTER)
    if i < len(steps)-1:
        ar = slide.shapes.add_connector(1, Inches(sx+sw), Inches(sy+0.55),
                                         Inches(sx+sw+0.2), Inches(sy+0.55))
        ar.line.color.rgb = LI_ACCENT; ar.line.width = Pt(2)

add_text(slide, "SSE Stream: GET /ai/task-stream/{task_id}  →  Real-time step progress to frontend",
         0.3, 3.5, 12.7, 0.35, size=11, color=LI_ACCENT, italic=True)

# Hiring Assistant Agent detail
add_rect(slide, 0.3, 3.95, 6.0, 2.8, fill=RGBColor(0x0A, 0x22, 0x38))
add_text(slide, "Hiring Assistant Agent", 0.45, 4.0, 5.7, 0.38, size=14, bold=True, color=LI_ACCENT)
hiring_steps = [
    "1. parse_resume()  → extract skills, experience",
    "2. compute_match_score()  → cosine similarity on skill vectors",
    "3. rank_candidates()  → top-N by weighted score",
    "4. Human-in-loop: recruiter reviews + approves/rejects",
    "5. approve_task()  → updates status → triggers analytics",
    "Approval Rate: tracked in MongoDB ai_db.tasks",
]
for j, s in enumerate(hiring_steps):
    add_text(slide, s, 0.5, 4.42+j*0.38, 5.7, 0.36, size=10, color=LI_WHITE)

# Career Coach Agent detail
add_rect(slide, 6.5, 3.95, 6.55, 2.8, fill=RGBColor(0x0A, 0x22, 0x38))
add_text(slide, "Career Coach Agent", 6.65, 4.0, 6.2, 0.38, size=14, bold=True, color=LI_ACCENT)
coach_steps = [
    "Input: member skills + target job role",
    "analyze_career_fit()  → skill gap analysis",
    "Outputs: skill_match_pct, overall_rating (1-10)",
    "matched_skills / missing_skills / bonus_skills",
    "headline_suggestion  → personalized tagline",
    "suggestions[]  → prioritized action items",
]
for j, s in enumerate(coach_steps):
    add_text(slide, s, 6.65, 4.42+j*0.38, 6.2, 0.36, size=10, color=LI_WHITE)

# Kafka topic legend
add_rect(slide, 0.3, 1.55, 12.7, 0.55, fill=RGBColor(0x1A, 0x30, 0x4A))
add_text(slide, "Kafka Topics: ", 0.45, 1.58, 1.4, 0.48, size=10, bold=True, color=LI_ACCENT)
topic_list = "job.created  ·  job.closed  ·  job.viewed  ·  job.saved  ·  application.submitted  ·  application.status_updated  ·  connection.requested  ·  connection.accepted  ·  message.sent  ·  member.created  ·  member.updated  ·  profile.viewed  ·  ai.requests  ·  ai.results"
add_text(slide, topic_list, 1.75, 1.58, 11.1, 0.48, size=9, color=LI_WHITE)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 – JMeter Performance: Throughput & Latency (Charts 1 & 2)
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, 13.33, 7.5, fill=LI_LIGHT)
slide_header(slide, "5 · Scalability & Performance  (Apache JMeter  ·  100 Users  ·  n=500)",
             "Job Search endpoint  ·  4 test combinations  ·  5 loops × 100 threads")
slide_footer(slide)

# B baseline callout box (too different in scale to chart alongside others)
add_rect(slide, 0.3, 1.3, 2.3, 3.6, fill=RGBColor(0xCC, 0x22, 0x00))
add_text(slide, "B\nBaseline\n(Cold Cache)", 0.35, 1.35, 2.2, 0.8, size=12, bold=True,
         color=LI_WHITE, align=PP_ALIGN.CENTER)
add_text(slide, "9.3\nTPS", 0.35, 2.15, 2.2, 0.9, size=28, bold=True,
         color=LI_WHITE, align=PP_ALIGN.CENTER)
add_text(slide, "1,565 ms\navg latency", 0.35, 3.1, 2.2, 0.65, size=16, bold=True,
         color=LI_WHITE, align=PP_ALIGN.CENTER)
add_text(slide, "No Redis\nNo cache", 0.35, 3.8, 2.2, 0.55, size=10,
         color=LI_WHITE, align=PP_ALIGN.CENTER, italic=True)

# Arrow
ar = slide.shapes.add_connector(1, Inches(2.6), Inches(3.1), Inches(2.85), Inches(3.1))
ar.line.color.rgb = LI_BLUE; ar.line.width = Pt(3)

# Chart 1: Throughput — B+S, B+S+K, B+S+K+X only (readable scale)
chart_data = ChartData()
chart_data.categories = ['B+S\n(+Redis)', 'B+S+K\n(+Kafka)', 'B+S+K+X\n(+Auth+FTS)']
chart_data.add_series('Transactions/sec', (101.9, 101.8, 101.6))

chart = slide.shapes.add_chart(
    XL_CHART_TYPE.COLUMN_CLUSTERED,
    Inches(2.85), Inches(1.3), Inches(4.8), Inches(3.6),
    chart_data
).chart
chart.has_title = True
chart.chart_title.text_frame.text = "Throughput (TPS) — with optimizations"
chart.chart_title.text_frame.paragraphs[0].font.size = Pt(11)
chart.chart_title.text_frame.paragraphs[0].font.bold = True
chart.chart_title.text_frame.paragraphs[0].font.color.rgb = LI_DARK
plot = chart.plots[0]
plot.has_data_labels = True
plot.data_labels.show_value = True
plot.data_labels.font.size = Pt(11)
plot.data_labels.font.bold = True
series = plot.series[0]
from pptx.oxml.ns import qn
from lxml import etree
for i, clr in enumerate(["057A55", "0A66C2", "E86A23"]):
    pt = series.points[i]
    pt.format.fill.solid()
    pt.format.fill.fore_color.rgb = RGBColor(int(clr[0:2],16), int(clr[2:4],16), int(clr[4:6],16))

# Chart 2: Avg Latency — B+S, B+S+K, B+S+K+X only (readable scale)
chart_data2 = ChartData()
chart_data2.categories = ['B+S\n(+Redis)', 'B+S+K\n(+Kafka)', 'B+S+K+X\n(+Auth+FTS)']
chart_data2.add_series('Avg Latency (ms)', (2, 3, 2))

chart2 = slide.shapes.add_chart(
    XL_CHART_TYPE.COLUMN_CLUSTERED,
    Inches(7.9), Inches(1.3), Inches(5.1), Inches(3.6),
    chart_data2
).chart
chart2.has_title = True
chart2.chart_title.text_frame.text = "Avg Latency (ms) — with optimizations"
chart2.chart_title.text_frame.paragraphs[0].font.size = Pt(11)
chart2.chart_title.text_frame.paragraphs[0].font.bold = True
chart2.chart_title.text_frame.paragraphs[0].font.color.rgb = LI_DARK
plot2 = chart2.plots[0]
plot2.has_data_labels = True
plot2.data_labels.show_value = True
plot2.data_labels.font.size = Pt(11)
plot2.data_labels.font.bold = True
series2 = plot2.series[0]
for i, clr in enumerate(["057A55", "0A66C2", "E86A23"]):
    pt = series2.points[i]
    pt.format.fill.solid()
    pt.format.fill.fore_color.rgb = RGBColor(int(clr[0:2],16), int(clr[2:4],16), int(clr[4:6],16))

# Key findings — JMeter real numbers
findings = [
    ("10.9×", "TPS gain: B (9.3) → B+S (101.9) with Redis"),
    ("99.9%", "Latency drop: 1,565ms → 2ms (cold→warm cache)"),
    ("<1ms", "Kafka overhead: B+S+K = 3ms vs B+S = 2ms"),
    ("0 errors", "Zero failed requests across all 4 JMeter runs"),
    ("p99 < 25ms", "P99 under 25ms for all cached combinations"),
]
add_rect(slide, 0.3, 5.1, 12.7, 0.38, fill=LI_DARK)
add_text(slide, "JMETER FINDINGS", 0.45, 5.12, 3, 0.34, size=11, bold=True, color=LI_WHITE)
add_text(slide, "Apache JMeter  ·  100 threads, 5s ramp, 5 loops  ·  500 samples/run  ·  B excluded from charts (scale: 1,565ms vs 2ms)",
         3.8, 5.13, 9.1, 0.32, size=8.5, color=LI_ACCENT, italic=True)
for i, (metric, desc) in enumerate(findings):
    col = i % 3; row = i // 3
    x = 0.3 + col * 4.35; y = 5.58 + row * 0.56
    add_rect(slide, x, y, 4.15, 0.5, fill=LI_WHITE)
    add_rect(slide, x, y, 0.06, 0.5, fill=LI_BLUE)
    add_text(slide, metric, x+0.12, y+0.04, 1.1, 0.42, size=13, bold=True, color=LI_BLUE)
    add_text(slide, desc,   x+1.2,  y+0.06, 2.85, 0.38, size=9, color=LI_DARK)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 – Agent-Related Metrics
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, 13.33, 7.5, fill=LI_LIGHT)
slide_header(slide, "6 · Agent-Related Metrics", "Recommendation Quality · Approval Rate · Skill Match Scores")
slide_footer(slide)

# Metric cards row 1
metrics = [
    ("Recommendation\nQuality Score", "8.4 / 10", "Avg AI rating across\nall hiring tasks", LI_BLUE),
    ("Approval Rate", "71%", "Recruiter approvals\n/ total tasks run", LI_GREEN),
    ("Avg Skill Match", "67%", "Candidate skills vs\njob requirements", LI_ORANGE),
    ("Tasks Completed", "41", "Applications processed\nby AI shortlister", LI_DARK),
]
for i, (title, val, sub, col) in enumerate(metrics):
    x = 0.3 + i * 3.2
    add_rect(slide, x, 1.3, 3.0, 1.95, fill=col)
    add_text(slide, title, x+0.1, 1.35, 2.8, 0.6, size=12, bold=True, color=LI_WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, val,   x+0.1, 1.9,  2.8, 0.75, size=30, bold=True, color=LI_WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, sub,   x+0.1, 2.65, 2.8, 0.52, size=10, color=LI_WHITE, align=PP_ALIGN.CENTER, italic=True)

# Caching / messaging policy details
add_rect(slide, 0.3, 3.45, 6.1, 3.42, fill=LI_WHITE)
add_rect(slide, 0.3, 3.45, 0.08, 3.42, fill=LI_BLUE)
add_text(slide, "Caching & Detection Policy", 0.5, 3.5, 5.7, 0.4, size=14, bold=True, color=LI_DARK)
cache_items = [
    "Redis: search results cached 60 s by sha256(params)",
    "Redis: single job cached 10 s by job_id",
    "Cache invalidated on: job update / job close",
    "Duplicate Kafka events: idempotency_key in processed_events",
    "at-least-once delivery guaranteed via MySQL dedup table",
    "Analytics consumer deduplicates by idempotency_key in MongoDB",
    "JWT expiry validation: custom HMAC-SHA256 in middleware",
    "Role-based access: member / recruiter enforced per-route",
]
for j, item in enumerate(cache_items):
    add_text(slide, f"• {item}", 0.5, 3.98+j*0.35, 5.8, 0.34, size=10, color=LI_DARK)

# Messaging flow details
add_rect(slide, 6.6, 3.45, 6.4, 3.42, fill=LI_WHITE)
add_rect(slide, 6.6, 3.45, 0.08, 3.42, fill=LI_GREEN)
add_text(slide, "Messaging Flow & Timeliness", 6.8, 3.5, 6.0, 0.4, size=14, bold=True, color=LI_DARK)
msg_items = [
    "message.sent → Kafka → SSE push to recipient",
    "Real-time delivery via Server-Sent Events (SSE)",
    "Thread.message_count incremented on each message",
    "application.submitted → job-service updates applicants_count",
    "job.viewed → async view counter increment (non-blocking)",
    "job.created → analytics ingestion for recruiter dashboard",
    "ai.requests → AI worker → ai.results → analytics store",
    "Career Coach: synchronous FastAPI (no Kafka needed)",
]
for j, item in enumerate(msg_items):
    add_text(slide, f"• {item}", 6.78, 3.98+j*0.35, 6.1, 0.34, size=10, color=LI_DARK)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 – Tech Stack Summary
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, 13.33, 7.5, fill=LI_DARK)
slide_footer(slide)

add_rect(slide, 0, 0, 13.33, 1.15, fill=LI_DARK)
add_rect(slide, 0, 1.15, 13.33, 0.06, fill=LI_BLUE)
add_text(slide, "Technology Stack & Scalability Techniques", 0.4, 0.12, 12, 0.7,
         size=26, bold=True, color=LI_WHITE)
add_text(slide, "How detection rules, caching policy, and messaging flows support timeliness and scalability",
         0.4, 0.75, 12, 0.38, size=13, color=LI_ACCENT)

stack = [
    ("Frontend",        "React 18 + Vite\nSSE real-time feeds\nVite proxy for all APIs", LI_BLUE),
    ("Profile API",     "FastAPI (Python)\nJWT auth (HS256)\nMySQL + Redis", RGBColor(0x0A,0x55,0xA0)),
    ("Job API",         "Node.js + Express\nMySQL full-text search\nRedis cache layer", LI_DARK),
    ("Application API", "Node.js + Express\nMultipart file upload\nKafka producer", RGBColor(0x02,0x4A,0x80)),
    ("Connection API",  "Node.js + Express\nMySQL\nKafka events", RGBColor(0x03,0x40,0x70)),
    ("Messaging API",   "Node.js + Express\nMongoDB + Mongoose\nSSE real-time", LI_GREEN),
    ("Analytics API",   "Node.js + Express\nMongoDB aggregation\nKafka consumer", RGBColor(0x03,0x60,0x40)),
    ("AI Service",      "FastAPI (Python)\nMongoDB task store\nKafka orchestration", LI_ORANGE),
]
for i, (name, desc, col) in enumerate(stack):
    col_idx = i % 4
    row_idx = i // 4
    x = 0.25 + col_idx * 3.25
    y = 1.4 + row_idx * 2.2
    add_rect(slide, x, y, 3.05, 2.0, fill=col)
    add_text(slide, name, x+0.12, y+0.1, 2.8, 0.45, size=13, bold=True, color=LI_WHITE)
    add_text(slide, desc, x+0.12, y+0.58, 2.8, 1.3, size=10.5, color=LI_WHITE)

add_text(slide, "All services containerized (Docker) · Orchestrated via docker-compose · Kafka via Confluent Platform · 78,985 job records seeded",
         0.3, 6.88, 12.7, 0.38, size=10, color=LI_ACCENT, italic=True, align=PP_ALIGN.CENTER)



# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 – Docker Containerization
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, 13.33, 7.5, fill=RGBColor(0x0D, 0x1B, 0x2A))
slide_header(slide, "Docker Containerization", "Every service ships as an independent Docker image · Orchestrated via docker-compose")
slide_footer(slide)

# Dockerfile pattern boxes
dockerfiles = [
    ("profile-service", "python:3.11-slim", "EXPOSE 8000", "uvicorn / run.py", LI_BLUE),
    ("job-service",     "node:20-alpine",   "EXPOSE 3002", "node src/server.js", LI_DARK),
    ("app-service",     "node:20-alpine",   "EXPOSE 5003", "node src/server.js", RGBColor(0x02,0x4A,0x80)),
    ("connection-svc",  "node:20-alpine",   "EXPOSE 3005", "node server.js", RGBColor(0x03,0x40,0x70)),
    ("messaging-svc",   "node:20-alpine",   "EXPOSE 3004", "node server.js", LI_GREEN),
    ("analytics-svc",   "node:20-alpine",   "EXPOSE 4000", "node src/server.js", RGBColor(0x03,0x60,0x40)),
    ("ai-service",      "python:3.11-slim", "EXPOSE 8005", "uvicorn app.main:app", LI_ORANGE),
    ("frontend",        "node:20 → nginx",  "EXPOSE 80",   "nginx (prod build)", RGBColor(0x60,0x20,0xA0)),
]
for i, (svc, base, expose, cmd, col) in enumerate(dockerfiles):
    ci = i % 4; ri = i // 4
    x = 0.3 + ci * 3.22; y = 1.35 + ri * 2.05
    add_rect(slide, x, y, 3.05, 1.88, fill=col)
    add_rect(slide, x, y, 3.05, 0.38, fill=RGBColor(0,0,0))
    add_text(slide, f"🐳  {svc}", x+0.1, y+0.04, 2.85, 0.32, size=11, bold=True, color=LI_WHITE)
    lines = [f"FROM {base}", f"WORKDIR /app", f"COPY & RUN install", expose, f"CMD {cmd}"]
    for j, line in enumerate(lines):
        add_text(slide, line, x+0.12, y+0.44+j*0.27, 2.82, 0.26, size=9, color=LI_WHITE)

# docker-compose highlights
add_rect(slide, 0.3, 5.6, 12.7, 1.45, fill=RGBColor(0x0A, 0x22, 0x38))
add_text(slide, "docker-compose.yml  —  Infrastructure Services", 0.45, 5.65, 8, 0.35, size=12, bold=True, color=LI_ACCENT)
infra = [
    ("Zookeeper", ":2181", "Kafka coordination"),
    ("Kafka (Confluent)", ":9092", "Event streaming broker"),
    ("Kafka UI", ":18088", "Topic & consumer monitoring"),
    ("MySQL 8.0", ":3308", "Transactional data"),
    ("MongoDB 7.0", ":27017", "Document store"),
    ("Redis 7", ":6379", "Cache layer"),
]
for i, (svc, port, desc) in enumerate(infra):
    col = i % 3; row = i // 3
    x = 0.4 + col * 4.3; y = 6.08 + row * 0.42
    add_text(slide, f"▸ {svc} {port}  —  {desc}", x, y, 4.1, 0.38, size=9.5, color=LI_WHITE)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 – Database Population Proof
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, 13.33, 7.5, fill=LI_LIGHT)
slide_header(slide, "Database Population  —  10,000+ Data Points", "Real production-scale seed data across all services")
slide_footer(slide)

# Big count cards
db_counts = [
    ("78,986",  "Job Postings",        "MySQL · data236",         LI_BLUE,  "Real job listings seeded\nfrom CSV dataset"),
    ("10,000+", "Members",             "MySQL · linkedin_sim",    LI_GREEN, "Synthetic member profiles\nwith skills & experience"),
    ("41",      "Applications",        "MySQL · application_db",  LI_DARK,  "Live user submissions\nduring testing"),
    ("4",       "Connections",         "MySQL · linkedin_conn",   RGBColor(0x03,0x40,0x70), "Member connection graph\n(pending/accepted)"),
    ("7",       "Messages",            "MongoDB · messaging",     LI_ORANGE,"Threaded conversations\nstored as documents"),
    ("14",      "Kafka Topics",        "Confluent Kafka",         RGBColor(0x60,0x10,0x80), "Event streams across\nall services"),
]
for i, (count, label, db, col, note) in enumerate(db_counts):
    ci = i % 3; ri = i // 3
    x = 0.4 + ci * 4.2; y = 1.35 + ri * 2.55
    add_rect(slide, x, y, 3.95, 2.3, fill=col)
    add_text(slide, count, x+0.15, y+0.1,  3.65, 1.0, size=46, bold=True, color=LI_WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, label, x+0.15, y+1.08, 3.65, 0.42, size=15, bold=True, color=LI_WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, db,    x+0.15, y+1.5,  3.65, 0.3,  size=9,  color=LI_WHITE, align=PP_ALIGN.CENTER, italic=True)
    add_text(slide, note,  x+0.15, y+1.82, 3.65, 0.4,  size=8.5,color=LI_WHITE, align=PP_ALIGN.CENTER)

add_text(slide, "Job postings seeded via CSV import (78,986 real listings) · Members generated with Faker · Full-text search index on title+description",
         0.4, 6.75, 12.5, 0.38, size=10, color=LI_GRAY, align=PP_ALIGN.CENTER, italic=True)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 12 – JMeter Charts 3 & 4: Concurrency Scaling & P90/P99 Latency
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, 13.33, 7.5, fill=LI_LIGHT)
slide_header(slide, "5b · JMeter  —  Concurrency Scaling & Percentile Latency",
             "Apache JMeter  ·  B+S (Redis cached)  ·  10 / 25 / 50 / 100 concurrent users")
slide_footer(slide)

# Chart 3: TPS vs concurrency (LINE — JMeter real)
chart_data3 = ChartData()
chart_data3.categories = ['10 users', '25 users', '50 users', '100 users']
chart_data3.add_series('Throughput (TPS)', (18.7, 43.9, 85.9, 170.6))

chart3 = slide.shapes.add_chart(
    XL_CHART_TYPE.LINE,
    Inches(0.4), Inches(1.3), Inches(5.9), Inches(3.5),
    chart_data3
).chart
chart3.has_title = True
chart3.chart_title.text_frame.text = "Throughput (TPS) vs Concurrency ↑"
chart3.chart_title.text_frame.paragraphs[0].font.size = Pt(12)
chart3.chart_title.text_frame.paragraphs[0].font.bold = True
chart3.chart_title.text_frame.paragraphs[0].font.color.rgb = LI_DARK
chart3.plots[0].has_data_labels = True
chart3.plots[0].data_labels.show_value = True
chart3.plots[0].data_labels.font.size = Pt(10)
chart3.plots[0].data_labels.font.bold = True

# Chart 4: P90 vs P99 — B+S, B+S+K, B+S+K+X only (B excluded: P99=18,513ms crushes scale)
# B callout box instead
add_rect(slide, 6.7, 1.3, 2.0, 3.5, fill=RGBColor(0xCC, 0x22, 0x00))
add_text(slide, "B Baseline\n(excluded\nfrom chart)", 6.75, 1.35, 1.9, 0.7, size=10, bold=True,
         color=LI_WHITE, align=PP_ALIGN.CENTER)
add_text(slide, "P90\n9,766ms", 6.75, 2.1, 1.9, 0.75, size=16, bold=True,
         color=LI_WHITE, align=PP_ALIGN.CENTER)
add_text(slide, "P99\n18,513ms", 6.75, 2.9, 1.9, 0.75, size=16, bold=True,
         color=LI_WHITE, align=PP_ALIGN.CENTER)
add_text(slide, "No Redis", 6.75, 3.7, 1.9, 0.35, size=9,
         color=LI_WHITE, align=PP_ALIGN.CENTER, italic=True)

chart_data4 = ChartData()
chart_data4.categories = ['B+S\n(+Redis)', 'B+S+K\n(+Kafka)', 'B+S+K+X\n(+Auth)']
chart_data4.add_series('P90 (ms)', (2, 3, 3))
chart_data4.add_series('P99 (ms)', (8, 24, 6))

chart4 = slide.shapes.add_chart(
    XL_CHART_TYPE.COLUMN_CLUSTERED,
    Inches(8.9), Inches(1.3), Inches(4.1), Inches(3.5),
    chart_data4
).chart
chart4.has_title = True
chart4.chart_title.text_frame.text = "P90 vs P99 Latency (ms) — optimized"
chart4.chart_title.text_frame.paragraphs[0].font.size = Pt(11)
chart4.chart_title.text_frame.paragraphs[0].font.bold = True
chart4.chart_title.text_frame.paragraphs[0].font.color.rgb = LI_DARK
plot4 = chart4.plots[0]
plot4.has_data_labels = True
plot4.data_labels.show_value = True
plot4.data_labels.font.size = Pt(10)
plot4.data_labels.font.bold = True
for i, clr in enumerate(["0A66C2", "E86A23"]):
    s = plot4.series[i]
    s.format.fill.solid()
    s.format.fill.fore_color.rgb = RGBColor(int(clr[0:2],16), int(clr[2:4],16), int(clr[4:6],16))

# JMeter results table
add_rect(slide, 0.3, 5.05, 12.7, 0.38, fill=LI_DARK)
add_text(slide, "JMETER SUMMARY TABLE", 0.45, 5.07, 4, 0.32, size=10, bold=True, color=LI_WHITE)
add_text(slide, "Tool: Apache JMeter  ·  All runs: 0 errors, 0 failed requests",
         5.5, 5.08, 7.3, 0.3, size=9, color=LI_ACCENT, italic=True)

headers = ["Test", "Threads", "Samples", "Avg (ms)", "P90 (ms)", "P99 (ms)", "TPS", "Errors"]
rows_data = [
    ["B (cold cache)",    "100", "500", "1,565", "9,766", "18,513", "9.3",  "0"],
    ["B+S (Redis warm)",  "100", "500", "2",     "2",     "8",      "101.9","0"],
    ["B+S+K (+Kafka)",    "100", "500", "3",     "3",     "24",     "101.8","0"],
    ["B+S+K+X (+Auth)",   "100", "500", "2",     "3",     "6",      "101.6","0"],
]
col_widths = [2.2, 0.85, 0.95, 0.95, 1.0, 1.0, 0.9, 0.8]
col_x = [0.3]
for w in col_widths[:-1]: col_x.append(col_x[-1]+w)

# Header row
for j, (hdr, cx, cw) in enumerate(zip(headers, col_x, col_widths)):
    add_rect(slide, cx, 5.52, cw-0.04, 0.34, fill=LI_BLUE)
    add_text(slide, hdr, cx+0.04, 5.54, cw-0.08, 0.3, size=8.5, bold=True,
             color=LI_WHITE, align=PP_ALIGN.CENTER)

for ri, row in enumerate(rows_data):
    bg = LI_WHITE if ri % 2 == 0 else RGBColor(0xE8, 0xF0, 0xFE)
    for j, (val, cx, cw) in enumerate(zip(row, col_x, col_widths)):
        add_rect(slide, cx, 5.9+ri*0.32, cw-0.04, 0.3, fill=bg)
        fc = LI_DARK if j > 0 else LI_BLUE
        add_text(slide, val, cx+0.04, 5.92+ri*0.32, cw-0.08, 0.26,
                 size=8.5, color=fc, bold=(j==0),
                 align=PP_ALIGN.CENTER if j > 0 else PP_ALIGN.LEFT)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 13 – Timeliness, Caching Policy & Detection Rules
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, 13.33, 7.5, fill=LI_LIGHT)
slide_header(slide, "Timeliness, Caching Policy & Detection Rules",
             "How our design supports scalability and real-time delivery")
slide_footer(slide)

sections = [
    ("🗄  SQL Caching Policy  (Redis)", LI_BLUE, [
        "Search results: cached by sha256(query params), TTL = 60 seconds",
        "Single job fetch: cached by job_id, TTL = 10 seconds",
        "Cache key format:  search:{sha256}  /  job:{uuid}",
        "Invalidation: job update or close flushes that job's key immediately",
        "Impact: 99.9% latency reduction (1,565ms → 2ms), 10.9× TPS gain (JMeter)",
        "Redis 7 Alpine — minimal memory footprint, sub-millisecond lookups",
    ]),
    ("📨  Messaging Flow & Timeliness", LI_GREEN, [
        "Kafka fire-and-forget: HTTP responds immediately, events publish async",
        "message.sent → SSE stream → real-time delivery to recipient browser",
        "application.submitted → job-service increments applicants_count async",
        "job.viewed → view counter updated without blocking the HTTP response",
        "ai.requests → worker processes → ai.results → analytics persisted",
        "All SSE streams use EventEmitter in-process bus for sub-100ms delivery",
    ]),
    ("🔍  Detection Rules & Idempotency", LI_ORANGE, [
        "Kafka at-least-once delivery — duplicate events WILL arrive",
        "MySQL processed_events table: idempotency_key PRIMARY KEY deduplicates",
        "MongoDB analytics: findOne({idempotency_key}) before insert",
        "UUID v4 idempotency_key in every event envelope — globally unique",
        "JWT expiry checked on every authenticated request (custom HMAC-SHA256)",
        "Role enforcement: member vs recruiter checked per-route, 403 on mismatch",
    ]),
]

for i, (title, col, bullets) in enumerate(sections):
    x = 0.3 + i * 4.35
    add_rect(slide, x, 1.3, 4.15, 5.65, fill=LI_WHITE)
    add_rect(slide, x, 1.3, 4.15, 0.42, fill=col)
    add_text(slide, title, x+0.1, 1.32, 3.95, 0.38, size=11, bold=True, color=LI_WHITE)
    for j, b in enumerate(bullets):
        add_text(slide, f"• {b}", x+0.12, 1.82+j*0.76, 3.9, 0.7, size=9.5, color=LI_DARK)


# ══════════════════════════════════════════════════════════════════════════════
# Save
# ══════════════════════════════════════════════════════════════════════════════
import os
os.makedirs("/Users/dipinjassal/sem_2/linkedin/presentation", exist_ok=True)
prs.save("/Users/dipinjassal/sem_2/linkedin/presentation/LinkedIn_Group6_Presentation.pptx")
print("Saved: /Users/dipinjassal/sem_2/linkedin/presentation/LinkedIn_Group6_Presentation.pptx")
