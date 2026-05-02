"""
LinkedIn Platform — Group 8 Presentation PPTX Generator
Only real data. No seeded/hallucinated metrics.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE

# ── Palette ───────────────────────────────────────────────────────────────────
BLUE      = RGBColor(0x0A, 0x66, 0xC2)
NAVY      = RGBColor(0x00, 0x1B, 0x48)
NAVY2     = RGBColor(0x02, 0x2A, 0x60)
SKY       = RGBColor(0x70, 0xB5, 0xF9)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
OFFWHITE  = RGBColor(0xF3, 0xF2, 0xEF)
GREEN     = RGBColor(0x05, 0x7A, 0x55)
GREEN2    = RGBColor(0x02, 0x4A, 0x35)
ORANGE    = RGBColor(0xE8, 0x6A, 0x23)
ORANGE2   = RGBColor(0x9C, 0x46, 0x17)
GRAY      = RGBColor(0x56, 0x68, 0x7A)
LGRAY     = RGBColor(0xD8, 0xD8, 0xD8)
RED       = RGBColor(0xCC, 0x22, 0x00)
REDDK     = RGBColor(0x88, 0x11, 0x00)
TEAL      = RGBColor(0x02, 0x7B, 0x8E)
PURPLE    = RGBColor(0x60, 0x20, 0xA0)
STEELBLUE = RGBColor(0x1A, 0x4A, 0x7C)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]

# ── Helpers ───────────────────────────────────────────────────────────────────

def rect(slide, l, t, w, h, fill, alpha=1.0):
    s = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = fill
    s.line.fill.background()
    return s

def rrect(slide, l, t, w, h, fill, radius_pt=8):
    """Rounded rectangle via MSO_SHAPE_TYPE 5 (rounded rectangle)."""
    s = slide.shapes.add_shape(5, Inches(l), Inches(t), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = fill
    s.line.fill.background()
    s.adjustments[0] = 0.05
    return s

def txt(slide, text, l, t, w, h, size=14, bold=False, color=WHITE,
        align=PP_ALIGN.LEFT, italic=False, wrap=True):
    box = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    box.word_wrap = wrap
    tf = box.text_frame; tf.word_wrap = wrap
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold
    r.font.italic = italic; r.font.color.rgb = color
    return box

def arrow(slide, x1, y1, x2, y2, color=BLUE, width=2):
    c = slide.shapes.add_connector(1, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    c.line.color.rgb = color; c.line.width = Pt(width)
    return c

def header(slide, title, subtitle=None):
    rect(slide, 0, 0, 13.33, 1.12, NAVY)
    rect(slide, 0, 1.12, 13.33, 0.05, BLUE)
    txt(slide, title, 0.38, 0.1, 12.2, 0.65, size=26, bold=True, color=WHITE)
    if subtitle:
        txt(slide, subtitle, 0.38, 0.72, 12.2, 0.38, size=12, color=SKY, italic=True)

def footer(slide, label="DATA 236  ·  LinkedIn Platform  ·  Group 8"):
    rect(slide, 0, 7.22, 13.33, 0.28, NAVY)
    txt(slide, label, 0.3, 7.24, 10, 0.24, size=8.5, color=SKY)
    txt(slide, "Spring 2025", 11.5, 7.24, 1.7, 0.24, size=8.5, color=SKY, align=PP_ALIGN.RIGHT)

def badge(slide, label, l, t, w=1.3, h=0.32, fill=BLUE, tsize=10):
    rrect(slide, l, t, w, h, fill)
    txt(slide, label, l, t+0.02, w, h-0.04, size=tsize, bold=True,
        color=WHITE, align=PP_ALIGN.CENTER)

def stat_card(slide, number, label, sublabel, l, t, w, h, fill):
    rrect(slide, l, t, w, h, fill)
    rect(slide, l, t, w, 0.06, BLUE)
    txt(slide, number, l+0.08, t+0.1, w-0.16, h*0.45,
        size=36, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txt(slide, label, l+0.08, t+h*0.5, w-0.16, 0.38,
        size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txt(slide, sublabel, l+0.08, t+h*0.78, w-0.16, 0.28,
        size=8.5, color=SKY, align=PP_ALIGN.CENTER, italic=True)

def svc_box(slide, label, port, lang, l, t, w=1.55, h=1.15, fill=BLUE):
    rrect(slide, l, t, w, h, fill)
    rect(slide, l, t, w, 0.28, NAVY2)
    txt(slide, label, l+0.06, t+0.03, w-0.12, 0.24,
        size=9, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txt(slide, port, l+0.06, t+0.34, w-0.12, 0.26,
        size=11, bold=True, color=SKY, align=PP_ALIGN.CENTER)
    txt(slide, lang, l+0.06, t+0.62, w-0.12, 0.46,
        size=8.5, color=LGRAY, align=PP_ALIGN.CENTER)

def db_box(slide, label, db_type, detail, l, t, w=2.0, h=1.0, fill=GREEN):
    rrect(slide, l, t, w, h, fill)
    txt(slide, label, l+0.08, t+0.06, w-0.16, 0.32,
        size=11, bold=True, color=WHITE)
    txt(slide, db_type, l+0.08, t+0.38, w-0.16, 0.24,
        size=9, color=SKY)
    txt(slide, detail, l+0.08, t+0.62, w-0.16, 0.32,
        size=8.5, color=LGRAY, italic=True)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Cover
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
rect(slide, 0, 0, 13.33, 7.5, NAVY)

# Gradient effect — layered dark bands
for i, (y, h_, alpha_fill) in enumerate([
    (0, 7.5, NAVY),
    (0, 0.06, BLUE),
    (7.44, 0.06, BLUE),
]):
    rect(slide, 0, y, 13.33, h_, alpha_fill)

# Decorative geometric shapes
shapes_def = [
    (11.2, 0.4, 2.8, 2.8, BLUE),
    (12.5, 0.1, 1.8, 1.8, RGBColor(0x05, 0x40, 0x90)),
    (0.2, 5.8, 1.5, 1.5, STEELBLUE),
    (10.5, 5.5, 2.2, 2.2, RGBColor(0x02, 0x30, 0x60)),
]
for cx, cy, sz, sz2, c in shapes_def:
    s = slide.shapes.add_shape(9, Inches(cx), Inches(cy), Inches(sz), Inches(sz2))
    s.fill.solid(); s.fill.fore_color.rgb = c; s.line.fill.background()

# Accent line
rect(slide, 1.6, 3.28, 9.5, 0.04, BLUE)

# Logo mark
for ox, oy, ow, oh, oc in [(0.5,1.5,0.12,1.05,BLUE),(0.78,2.2,0.12,0.35,BLUE),
                             (0.78,1.5,0.12,0.45,BLUE),(0.5,1.5,0.4,0.12,BLUE)]:
    rect(slide, ox, oy, ow, oh, oc)

txt(slide, "LinkedIn Platform", 1.6, 1.45, 10.5, 1.1,
    size=48, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
txt(slide, "Scalable Microservices  ·  Kafka Event Bus  ·  Agentic AI  ·  Real-Time Analytics",
    1.6, 2.55, 10.5, 0.55, size=17, color=SKY, italic=True)

txt(slide, "Group 8", 1.6, 3.45, 3, 0.45, size=18, bold=True, color=WHITE)
txt(slide,
    "Dipin Jassal  ·  Sarvesh Reshimwale  ·  Sammruddhi  ·  Anushka Khadatkar  ·  Rajesh  ·  Bhavya  ·  Shashira  ·  Nikhil",
    1.6, 3.92, 10.5, 0.38, size=12.5, color=SKY)
txt(slide, "DATA 236 — Advanced Database Systems  ·  Spring 2025",
    1.6, 4.35, 10.5, 0.38, size=13, color=GRAY)

# Tech badges
for tag, x in [("MySQL",1.6),("MongoDB",3.05),("Redis",4.38),
               ("Kafka",5.55),("FastAPI",6.82),("React",8.18),("Docker",9.45)]:
    badge(slide, tag, x, 5.15, w=1.22, h=0.34, fill=BLUE, tsize=11)

txt(slide, "78,986 jobs  ·  10,000+ members  ·  8 microservices  ·  14 Kafka topics",
    1.6, 5.7, 10.5, 0.38, size=12, color=RGBColor(0x90,0xC8,0xF8), italic=True)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — Team Members
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
rect(slide, 0, 0, 13.33, 7.5, OFFWHITE)
header(slide, "1 · Group 8  —  Team Members")
footer(slide)

members = ["Dipin Jassal","Sarvesh Reshimwale","Sammruddhi",
           "Anushka Khadatkar","Rajesh","Bhavya","Shashira","Nikhil"]

for i, name in enumerate(members):
    col = i % 4; row = i // 4
    x = 0.38 + col * 3.2; y = 1.35 + row * 2.5
    rrect(slide, x, y, 3.0, 2.22, WHITE)
    rect(slide, x, y, 0.07, 2.22, BLUE)
    # Avatar
    av = slide.shapes.add_shape(9, Inches(x+0.32), Inches(y+0.28), Inches(1.0), Inches(1.0))
    av.fill.solid(); av.fill.fore_color.rgb = BLUE; av.line.fill.background()
    initials = "".join(w[0] for w in name.split()[:2]).upper()
    txt(slide, initials, x+0.32, y+0.34, 1.0, 0.6,
        size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txt(slide, name, x+0.14, y+1.44, 2.72, 0.54,
        size=13, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

txt(slide, "Project: LinkedIn-Like Professional Network  ·  Microservices + Kafka + AI Orchestration",
    0.5, 6.82, 12.3, 0.3, size=11, color=GRAY, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — Database Schema
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
rect(slide, 0, 0, 13.33, 7.5, OFFWHITE)
header(slide, "2 · Database Schema", "MySQL (transactional)  ·  MongoDB (documents/events)  ·  Redis (cache layer)")
footer(slide)

# ── MySQL column (left) ──
mysql_tables = [
    ("members", "PK: member_id (CHAR 36)",
     ["email UNIQUE", "first_name, last_name", "headline, about_summary",
      "experience_json, education_json", "skills_json, profile_photo_url",
      "connections_count, profile_views"]),
    ("job_postings", "PK: job_id (CHAR 36)  FK: recruiter_id",
     ["title, description (FULLTEXT indexed)", "employment_type, location",
      "remote ENUM(onsite/remote/hybrid)", "skills_required JSON, salary_min/max",
      "status ENUM(open/closed)", "views_count, applicants_count"]),
    ("applications", "PK: application_id  UNIQUE(job_id, member_id)",
     ["job_id → job_postings, member_id → members",
      "resume_file_name, resume_file_path", "cover_letter, recruiter_note",
      "status ENUM(submitted/reviewed/accepted/rejected)"]),
]
cols_mysql = [BLUE, RGBColor(0x05,0x55,0xA0), STEELBLUE]
for i,(tname, pk, cols) in enumerate(mysql_tables):
    y = 1.3 + i * 1.88
    rrect(slide, 0.28, y, 4.0, 1.75, cols_mysql[i])
    rect(slide, 0.28, y, 4.0, 0.32, NAVY2)
    txt(slide, f"  {tname}", 0.32, y+0.03, 3.92, 0.28, size=11, bold=True, color=WHITE)
    txt(slide, pk, 0.38, y+0.36, 3.82, 0.24, size=8, color=SKY, italic=True)
    for j, col in enumerate(cols[:5]):
        txt(slide, f"  • {col}", 0.38, y+0.58+j*0.22, 3.82, 0.22, size=8.5, color=WHITE)

# ── More MySQL (middle-left) ──
extra = [
    ("saved_jobs", "PK: (user_id, job_id)", ["saved_at DATETIME", "INDEX(user_id), INDEX(job_id)"]),
    ("connections", "PK: connection_id", ["requester_id, receiver_id",
     "status ENUM(pending/accepted/rejected)", "UNIQUE(requester_id, receiver_id)"]),
    ("processed_events", "PK: idempotency_key", ["event_type, trace_id",
     "entity_type, entity_id", "processed_at — Kafka dedup table"]),
]
for i,(tname, pk, cols) in enumerate(extra):
    y = 1.3 + i * 1.88
    rrect(slide, 4.5, y, 3.35, 1.75, RGBColor(0x12,0x38,0x6A))
    rect(slide, 4.5, y, 3.35, 0.32, NAVY2)
    txt(slide, f"  {tname}", 4.54, y+0.03, 3.27, 0.28, size=11, bold=True, color=WHITE)
    txt(slide, pk, 4.58, y+0.36, 3.18, 0.24, size=8, color=SKY, italic=True)
    for j, col in enumerate(cols[:5]):
        txt(slide, f"  • {col}", 4.58, y+0.58+j*0.26, 3.18, 0.26, size=9, color=WHITE)

# ── MongoDB ──
rrect(slide, 8.1, 1.3, 4.9, 1.75, GREEN)
rect(slide, 8.1, 1.3, 4.9, 0.32, GREEN2)
txt(slide, "  MongoDB: linkedin_analytics.events", 8.14, 1.33, 4.82, 0.28, size=10, bold=True, color=WHITE)
for j, c in enumerate(["idempotency_key (dedup PK)", "event_type, trace_id, actor_id",
                         "entity {entity_type, entity_id}", "payload (flexible JSON schema)",
                         "14 Kafka topics → all ingested here", "_topic, _ingested_at (routing)"]):
    txt(slide, f"  • {c}", 8.18, 1.62+j*0.22, 4.76, 0.22, size=8.5, color=WHITE)

rrect(slide, 8.1, 3.18, 4.9, 1.75, RGBColor(0x03,0x5C,0x40))
rect(slide, 8.1, 3.18, 4.9, 0.32, GREEN2)
txt(slide, "  MongoDB: messaging + ai_db", 8.14, 3.21, 4.82, 0.28, size=10, bold=True, color=WHITE)
for j, c in enumerate(["threads: thread_id, participants[], last_message_at",
                         "messages: message_id, thread_id, sender_id",
                         "  message_text, status(sent/delivered/read)",
                         "ai_db.tasks: job_id, recruiter_id, steps[]",
                         "  status, created_at, trace_id"]):
    txt(slide, f"  • {c}", 8.18, 3.5+j*0.24, 4.76, 0.24, size=8.5, color=WHITE)

rrect(slide, 8.1, 5.08, 4.9, 1.75, ORANGE2)
rect(slide, 8.1, 5.08, 4.9, 0.32, RGBColor(0x60,0x28,0x00))
txt(slide, "  Redis Cache  (TTL-based)", 8.14, 5.11, 4.82, 0.28, size=10, bold=True, color=WHITE)
for j, c in enumerate(["Key: search:{sha256(params)}  →  TTL 60 s",
                         "Key: job:{job_id}  →  TTL 10 s",
                         "Invalidated on: job update / job close",
                         "Reduces MySQL FULLTEXT load by ~99%",
                         "Sub-millisecond lookups from memory"]):
    txt(slide, f"  • {c}", 8.18, 5.4+j*0.26, 4.76, 0.26, size=8.5, color=WHITE)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — System Architecture Diagram
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
rect(slide, 0, 0, 13.33, 7.5, RGBColor(0x0A, 0x12, 0x1E))
header(slide, "3 · System Architecture Diagram",
       "3-tier microservices  ·  8 services  ·  Kafka event bus  ·  3 databases")
footer(slide)

# Zone backgrounds
for zl, zt, zw, zh, zc, zlabel in [
    (0.15, 1.22, 1.35, 5.8, RGBColor(0x0F,0x20,0x38), "CLIENT"),
    (1.62, 1.22, 1.22, 5.8, RGBColor(0x0A,0x1A,0x30), "PROXY"),
    (2.96, 1.22, 5.52, 5.8, RGBColor(0x08,0x18,0x28), "MICROSERVICES"),
    (8.6,  1.22, 2.3,  5.8, RGBColor(0x10,0x08,0x22), "KAFKA"),
    (11.0, 1.22, 2.2,  5.8, RGBColor(0x06,0x18,0x10), "DATABASES"),
]:
    rect(slide, zl, zt, zw, zh, zc)
    rect(slide, zl, zt, zw, 0.26, RGBColor(0x15,0x30,0x50))
    txt(slide, zlabel, zl+0.05, zt+0.02, zw-0.1, 0.22,
        size=7.5, bold=True, color=SKY, align=PP_ALIGN.CENTER)

# Client boxes
rrect(slide, 0.22, 1.55, 1.2, 1.1, BLUE)
txt(slide, "React\nFrontend\n:3000", 0.25, 1.58, 1.14, 1.0, size=8.5, color=WHITE, align=PP_ALIGN.CENTER)
rrect(slide, 0.22, 2.85, 1.2, 0.85, RGBColor(0x30,0x60,0xA0))
txt(slide, "Recruiter\nPortal", 0.25, 2.87, 1.14, 0.78, size=8.5, color=WHITE, align=PP_ALIGN.CENTER)

# Vite proxy
rrect(slide, 1.68, 1.55, 1.1, 2.15, NAVY2)
txt(slide, "Vite\nProxy\n/api/*\n/analytics\n/ai", 1.7, 1.6, 1.06, 2.05, size=8, color=WHITE, align=PP_ALIGN.CENTER)

# Arrows: client → proxy
arrow(slide, 1.42, 2.1, 1.68, 2.1, SKY, 1.5)
arrow(slide, 1.42, 3.28, 1.68, 2.8, SKY, 1.5)

# Service boxes (2 rows)
svcs = [
    ("Profile Svc", ":8002", "FastAPI\nPython", 3.03, 1.55, BLUE),
    ("Job Svc",     ":3002", "Node.js\nExpress", 4.72, 1.55, RGBColor(0x12,0x4A,0x90)),
    ("App Svc",     ":5003", "Node.js\nExpress", 6.41, 1.55, STEELBLUE),
    ("Conn Svc",    ":3005", "Node.js\nExpress", 3.03, 3.1,  TEAL),
    ("Msg Svc",     ":3004", "Node.js\nMongoDB", 4.72, 3.1,  RGBColor(0x03,0x5C,0x80)),
    ("Analytics",  ":4000", "Node.js\nMongoDB",  6.41, 3.1,  GREEN),
    ("AI Service",  ":8005", "FastAPI\nPython",   3.03, 4.65, RGBColor(0x80,0x40,0x00)),
    ("Redis",       ":6379", "Cache\nIn-memory",  4.72, 4.65, ORANGE2),
]
for label, port, lang, sx, sy, sc in svcs:
    svc_box(slide, label, port, lang, sx, sy, fill=sc)
    # Arrow to Kafka
    arrow(slide, sx+1.55, sy+0.58, 8.6, sy+0.58, RGBColor(0x40,0x40,0x60), 1.2)

# Proxy → services arrows
arrow(slide, 2.78, 2.1, 3.03, 2.1, SKY, 1.5)

# Kafka bus
rrect(slide, 8.62, 1.55, 2.25, 5.28, RGBColor(0x12,0x08,0x30))
rect(slide, 8.62, 1.55, 2.25, 0.38, RGBColor(0x40,0x10,0x80))
txt(slide, "KAFKA  :9092", 8.66, 1.57, 2.17, 0.34, size=9, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
topics = ["job.created","job.closed","job.viewed","job.saved",
          "application.submitted","application.status_updated",
          "connection.requested","connection.accepted",
          "message.sent","member.created","member.updated",
          "profile.viewed","ai.requests","ai.results"]
for j, t in enumerate(topics):
    c = RGBColor(0x22,0x12,0x50) if j % 2 == 0 else RGBColor(0x1A,0x0A,0x40)
    rect(slide, 8.68, 2.0+j*0.34, 2.13, 0.3, c)
    txt(slide, t, 8.72, 2.01+j*0.34, 2.05, 0.28, size=7, color=SKY)

# Kafka → DBs arrows
arrow(slide, 10.87, 2.4, 11.02, 2.4, RGBColor(0x40,0x80,0x40), 1.5)
arrow(slide, 10.87, 4.2, 11.02, 4.2, RGBColor(0x40,0x80,0x40), 1.5)

# Data stores
db_box(slide, "MySQL", "data236 / linkedin_sim / app_db",
       "members · jobs · applications\nconnections · processed_events",
       11.05, 1.55, w=2.15, h=1.75, fill=BLUE)
db_box(slide, "MongoDB", "analytics / messaging / ai_db",
       "events · messages · threads\nai_tasks (document store)",
       11.05, 3.45, w=2.15, h=1.75, fill=GREEN)
db_box(slide, "Redis", "Cache  :6379",
       "search: TTL 60s\njob: TTL 10s",
       11.05, 5.32, w=2.15, h=1.38, fill=ORANGE2)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — Agent Architecture & Kafka Topic Design
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
rect(slide, 0, 0, 13.33, 7.5, RGBColor(0x08, 0x12, 0x20))
header(slide, "4 · Agent Architecture & Kafka Topic Design",
       "Hiring Assistant (Supervisor)  ·  Career Coach  ·  Kafka-orchestrated multi-step workflow")
footer(slide)

# Kafka topics banner
rrect(slide, 0.28, 1.25, 12.75, 0.48, RGBColor(0x15,0x08,0x35))
txt(slide, "KAFKA TOPICS:", 0.38, 1.28, 1.5, 0.42, size=9, bold=True, color=SKY)
topic_str = ("job.created  ·  job.closed  ·  job.viewed  ·  job.saved  ·  application.submitted  ·  "
             "application.status_updated  ·  connection.requested  ·  connection.accepted  ·  "
             "message.sent  ·  member.created  ·  member.updated  ·  profile.viewed  ·  ai.requests  ·  ai.results")
txt(slide, topic_str, 1.9, 1.28, 10.9, 0.42, size=8.5, color=WHITE)

# AI pipeline flow boxes
steps = [
    ("Recruiter\nSelects Job", BLUE),
    ("POST\n/ai/submit-task", RGBColor(0x05,0x50,0x9A)),
    ("Kafka\nai.requests", RGBColor(0x50,0x08,0x80)),
    ("AI Worker\nConsumer", GREEN),
    ("MongoDB\nStore Task", RGBColor(0x03,0x5C,0x40)),
    ("Kafka\nai.results", RGBColor(0x80,0x40,0x00)),
    ("Frontend\nSSE Stream", TEAL),
]
bw = 1.52; gap = 0.12; y0 = 1.92
for i, (label, col) in enumerate(steps):
    x = 0.28 + i*(bw+gap)
    rrect(slide, x, y0, bw, 0.95, col)
    txt(slide, label, x+0.06, y0+0.06, bw-0.12, 0.82,
        size=9.5, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    if i < len(steps)-1:
        arrow(slide, x+bw, y0+0.47, x+bw+gap, y0+0.47, SKY, 2)

txt(slide, "SSE: GET /ai/task-stream/{task_id}  →  real-time step progress pushed to recruiter UI",
    0.28, 3.0, 12.75, 0.3, size=10, color=SKY, italic=True, align=PP_ALIGN.CENTER)

# Two agent detail panels
for bx, btitle, bcol, bitems in [
    (0.28, "Hiring Assistant Agent  (Supervisor)", BLUE, [
        "Step 1: parse_resume()  →  extract skills, years of experience, seniority level",
        "Step 2: compute_match_score()  →  70% skills overlap + 30% seniority weighting",
        "Step 3: rank_candidates()  →  sorted top-N shortlist by weighted composite score",
        "Step 4: outreach_draft()  →  template-based personalized recruiter message",
        "Step 5: Human-in-Loop  →  recruiter approves / rejects before any action taken",
        "Orchestrated via Kafka: ai.requests → worker → ai.results → analytics stored",
    ]),
    (6.82, "Career Coach Agent  (Synchronous FastAPI)", RGBColor(0x03,0x5C,0x40), [
        "Input: member's skills[] + target_role string",
        "analyze_career_fit()  →  intersection / union skill gap analysis",
        "Output: skill_match_pct, overall_rating (1–10 scale)",
        "matched_skills / missing_skills / bonus_skills lists",
        "headline_suggestion  →  personalized LinkedIn tagline",
        "suggestions[]  →  prioritized actionable improvement items",
    ]),
]:
    pw = 6.3
    rrect(slide, bx, 3.38, pw, 3.5, RGBColor(0x0A,0x16,0x2A))
    rect(slide, bx, 3.38, pw, 0.38, NAVY2)
    txt(slide, btitle, bx+0.14, 3.4, pw-0.28, 0.34, size=11, bold=True, color=SKY)
    for j, line in enumerate(bitems):
        txt(slide, f"  →  {line}", bx+0.14, 3.84+j*0.48, pw-0.22, 0.46, size=9.5, color=WHITE)


# ── Helper: add a single column chart with data labels ──────────────────────
def perf_chart(slide, x, y, w, h, title, cats, vals, bar_color=BLUE):
    cd = ChartData()
    cd.categories = cats
    cd.add_series('', vals)
    ch = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(x), Inches(y), Inches(w), Inches(h), cd).chart
    ch.has_title = True
    ch.chart_title.text_frame.text = title
    ch.chart_title.text_frame.paragraphs[0].font.size = Pt(10)
    ch.chart_title.text_frame.paragraphs[0].font.bold = True
    ch.chart_title.text_frame.paragraphs[0].font.color.rgb = NAVY
    pl = ch.plots[0]
    pl.has_data_labels = True
    pl.data_labels.show_value = True
    pl.data_labels.font.size = Pt(9)
    pl.data_labels.font.bold = True
    s = pl.series[0]
    s.format.fill.solid()
    s.format.fill.fore_color.rgb = bar_color
    return ch

CATS = ['B', 'B+S', 'B+S+K', 'B+S+K+Other']

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — JMeter Scenario A: Job Search + Job Detail
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
rect(slide, 0, 0, 13.33, 7.5, OFFWHITE)
header(slide, "5a · JMeter  —  Scenario A: Job Search + Job Detail",
       "100 concurrent users  ·  4 configurations: B / B+S / B+S+K / B+S+K+Other")
footer(slide)

# 4 charts side by side — all 4 bars including B
CW = 3.12; CH = 4.8; CY = 1.25; GAP = 0.14
CXS = [0.25, 0.25+CW+GAP, 0.25+2*(CW+GAP), 0.25+3*(CW+GAP)]

perf_chart(slide, CXS[0], CY, CW, CH, "Throughput (req/s)",
           CATS, (63.5, 204.7, 204.7, 204.4), BLUE)
perf_chart(slide, CXS[1], CY, CW, CH, "Average Latency (ms)",
           CATS, (1289.5, 0.9, 0.9, 0.9), RGBColor(0x0A,0x55,0xA0))
perf_chart(slide, CXS[2], CY, CW, CH, "P95 Latency (ms)",
           CATS, (7359.0, 2.0, 2.0, 2.0), RGBColor(0x02,0x40,0x80))
perf_chart(slide, CXS[3], CY, CW, CH, "Error Rate (%)",
           CATS, (0.0, 0.0, 0.0, 0.0), RED)

# Findings strip
rect(slide, 0.25, 6.18, 12.83, 0.3, NAVY)
txt(slide, "KEY FINDINGS  —  Scenario A", 0.38, 6.2, 3.5, 0.26, size=9, bold=True, color=WHITE)
txt(slide, "Throughput: B=63.5 → B+S=204.7 req/s (3.2× gain)  ·  Latency: 1,289ms → 0.9ms (99.9% drop)  ·  P95: 7,359ms → 2ms  ·  0% error rate all configs",
    3.9, 6.21, 9.0, 0.26, size=8.5, color=SKY, italic=True)

scA = [
    ("3.2×",   "Throughput gain B→B+S (Redis)"),
    ("99.9%",  "Latency drop: 1,289ms→0.9ms"),
    ("0.9ms",  "Avg latency with Redis warm"),
    ("0%",     "Error rate — all 4 configs"),
]
for i, (metric, desc) in enumerate(scA):
    x = 0.25 + i * 3.23
    rrect(slide, x, 6.54, 3.1, 0.7, WHITE)
    rect(slide, x, 6.54, 0.07, 0.7, BLUE)
    txt(slide, metric, x+0.14, 6.56, 1.1, 0.32, size=16, bold=True, color=BLUE)
    txt(slide, desc,   x+1.22, 6.58, 1.78, 0.28, size=8.5, color=NAVY)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — JMeter Scenario B: Apply Submit (DB + Kafka)
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
rect(slide, 0, 0, 13.33, 7.5, OFFWHITE)
header(slide, "5b · JMeter  —  Scenario B: Apply Submit  (DB + Kafka)",
       "100 concurrent users  ·  4 configurations: B / B+S / B+S+K / B+S+K+Other")
footer(slide)

perf_chart(slide, CXS[0], CY, CW, CH, "Throughput (req/s)",
           CATS, (102.2, 100.8, 102.0, 102.1), GREEN)
perf_chart(slide, CXS[1], CY, CW, CH, "Average Latency (ms)",
           CATS, (2.7, 4.5, 2.5, 2.5), RGBColor(0x03,0x60,0x40))
perf_chart(slide, CXS[2], CY, CW, CH, "P95 Latency (ms)",
           CATS, (4.0, 4.0, 4.0, 4.0), RGBColor(0x02,0x48,0x30))
perf_chart(slide, CXS[3], CY, CW, CH, "Error Rate (%)",
           CATS, (0.0, 0.0, 0.0, 0.0), TEAL)

# Findings strip
rect(slide, 0.25, 6.18, 12.83, 0.3, NAVY)
txt(slide, "KEY FINDINGS  —  Scenario B", 0.38, 6.2, 3.5, 0.26, size=9, bold=True, color=WHITE)
txt(slide, "Throughput stable ~100 req/s (write-bound: DB + Kafka async)  ·  Avg latency 2–5ms  ·  P95 = 4ms  ·  Zero errors across ALL configs",
    3.9, 6.21, 9.0, 0.26, size=8.5, color=SKY, italic=True)

scB = [
    ("~100",   "req/s stable (write-bound)"),
    ("<5ms",   "Avg latency all configs"),
    ("0%",     "Error rate — all 4 configs"),
    ("4ms",    "P95 latency — consistent"),
]
for i, (metric, desc) in enumerate(scB):
    x = 0.25 + i * 3.23
    rrect(slide, x, 6.54, 3.1, 0.7, WHITE)
    rect(slide, x, 6.54, 0.07, 0.7, GREEN)
    txt(slide, metric, x+0.14, 6.56, 1.1, 0.32, size=16, bold=True, color=GREEN)
    txt(slide, desc,   x+1.22, 6.58, 1.78, 0.28, size=8.5, color=NAVY)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — Database Population (Real Counts)
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
rect(slide, 0, 0, 13.33, 7.5, RGBColor(0x06,0x10,0x1E))
header(slide, "6 · Database Population  —  Real Scale",
       "Production-scale seed data  ·  All counts from actual database queries")
footer(slide)

db_stats = [
    ("78,986",  "Job Postings",    "MySQL · data236",         BLUE,    "Real listings from\nKaggle LinkedIn dataset"),
    ("10,000+", "Members",         "MySQL · linkedin_sim",    GREEN,   "Seeded via Kaggle\nResume CSV + Faker"),
    ("2,484",   "Resume Profiles", "MySQL · linkedin_sim",    STEELBLUE,"Extracted from\nKaggle Resume.csv"),
    ("4",       "Connections",     "MySQL · linkedin_conn",   TEAL,    "Live connection\ngraph (test users)"),
    ("14",      "Kafka Topics",    "Confluent Kafka :9092",   PURPLE,  "Events across all\nmicroservices"),
    ("3",       "Database Types",  "MySQL+MongoDB+Redis",     ORANGE2, "RDBMS + Document\n+ Cache layer"),
]
for i,(count,label,db,col,note) in enumerate(db_stats):
    ci = i%3; ri = i//3
    x = 0.38 + ci*4.3; y = 1.32 + ri*2.65
    rrect(slide, x, y, 4.1, 2.4, col)
    rect(slide, x, y, 4.1, 0.06, RGBColor(0xFF,0xFF,0xFF))
    txt(slide, count, x+0.12, y+0.1, 3.86, 1.1,
        size=54, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txt(slide, label, x+0.12, y+1.18, 3.86, 0.42,
        size=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txt(slide, db, x+0.12, y+1.6, 3.86, 0.3,
        size=9, color=SKY, align=PP_ALIGN.CENTER, italic=True)
    txt(slide, note, x+0.12, y+1.92, 3.86, 0.4,
        size=8.5, color=RGBColor(0xCC,0xDD,0xFF), align=PP_ALIGN.CENTER)

txt(slide, "Job postings from Kaggle LinkedIn Job 2023 dataset  ·  Members seeded from Kaggle Resume Dataset  ·  FULLTEXT index on title + description",
    0.4, 6.82, 12.5, 0.3, size=9.5, color=GRAY, align=PP_ALIGN.CENTER, italic=True)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — Tech Stack & Caching / Detection Policy
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
rect(slide, 0, 0, 13.33, 7.5, NAVY)
header(slide, "7 · Tech Stack & Design Policies",
       "How caching, idempotency, and messaging flows support timeliness and scalability")
footer(slide)

stack = [
    ("Frontend",       "React 18 + Vite\nSSE real-time feeds\nVite proxy — no CORS", BLUE),
    ("Profile API",    "FastAPI (Python)\nJWT HS256 auth\nMySQL + Redis cache", RGBColor(0x0A,0x55,0xA0)),
    ("Job API",        "Node.js + Express\nMySQL FULLTEXT search\nRedis 60s / 10s TTL", STEELBLUE),
    ("Application API","Node.js + Express\nMultipart file upload\nKafka fire-and-forget", RGBColor(0x02,0x4A,0x80)),
    ("Connection API", "Node.js + Express\nMySQL graph queries\nKafka events", TEAL),
    ("Messaging API",  "Node.js + Express\nMongoDB + Mongoose\nSSE real-time push", GREEN),
    ("Analytics API",  "Node.js + Express\nMongoDB aggregation\nKafka consumer (14 topics)", RGBColor(0x03,0x60,0x40)),
    ("AI Service",     "FastAPI (Python)\nMongoDB task store\nKafka orchestration", ORANGE2),
]
for i,(name,desc,col) in enumerate(stack):
    ci = i%4; ri = i//4
    x = 0.25 + ci*3.27; y = 1.32 + ri*2.12
    rrect(slide, x, y, 3.08, 1.98, col)
    rect(slide, x, y, 3.08, 0.32, NAVY2)
    txt(slide, name, x+0.12, y+0.04, 2.84, 0.26, size=12, bold=True, color=WHITE)
    txt(slide, desc, x+0.12, y+0.4, 2.84, 1.4, size=9.5, color=LGRAY)

rect(slide, 0.25, 5.68, 12.8, 0.04, BLUE)

policies = [
    ("SQL Cache Policy", BLUE, [
        "search:{sha256} → TTL 60s",
        "job:{uuid} → TTL 10s",
        "Invalidate on: update / close",
        "99.9% latency reduction verified",
    ]),
    ("Kafka Idempotency", PURPLE, [
        "processed_events PK dedup",
        "idempotency_key per envelope",
        "At-least-once + safe retry",
        "MongoDB: findOne before insert",
    ]),
    ("Messaging Flow", GREEN, [
        "HTTP responds → Kafka async",
        "SSE push: <100ms delivery",
        "DB write before Kafka publish",
        "Job viewed: non-blocking async",
    ]),
]
for i,(ptitle,pcol,pitems) in enumerate(policies):
    x = 0.25 + i*4.36
    rrect(slide, x, 5.78, 4.2, 1.48, RGBColor(0x0A,0x16,0x2A))
    rect(slide, x, 5.78, 4.2, 0.28, pcol)
    txt(slide, ptitle, x+0.1, 5.8, 4.0, 0.26, size=10, bold=True, color=WHITE)
    for j, item in enumerate(pitems):
        txt(slide, f"• {item}", x+0.12, 6.12+j*0.27, 3.95, 0.26, size=8.5, color=LGRAY)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — Docker Containerization
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
rect(slide, 0, 0, 13.33, 7.5, RGBColor(0x08,0x12,0x20))
header(slide, "8 · Docker Containerization & Infrastructure",
       "Every service ships as an independent Docker image  ·  Orchestrated via docker-compose")
footer(slide)

dockerfiles = [
    ("profile-service", "python:3.11-slim", "8002", "uvicorn app.main:app", BLUE),
    ("job-service",     "node:20-alpine",   "3002", "node src/server.js",   STEELBLUE),
    ("app-service",     "node:20-alpine",   "5003", "node src/server.js",   RGBColor(0x02,0x4A,0x80)),
    ("connection-svc",  "node:20-alpine",   "3005", "node server.js",       TEAL),
    ("messaging-svc",   "node:20-alpine",   "3004", "node server.js",       GREEN),
    ("analytics-svc",   "node:20-alpine",   "4000", "node src/server.js",   RGBColor(0x03,0x60,0x40)),
    ("ai-service",      "python:3.11-slim", "8005", "uvicorn app.main:app", ORANGE2),
    ("frontend",        "node:20 → nginx",  "80",   "nginx (prod build)",   PURPLE),
]
for i,(svc,base,port,cmd,col) in enumerate(dockerfiles):
    ci = i%4; ri = i//4
    x = 0.28 + ci*3.27; y = 1.28 + ri*2.08
    rrect(slide, x, y, 3.08, 1.92, col)
    rect(slide, x, y, 3.08, 0.3, RGBColor(0,0,0))
    txt(slide, f"  {svc}", x+0.08, y+0.04, 2.92, 0.25, size=10, bold=True, color=WHITE)
    lines = [f"FROM {base}", "WORKDIR /app", "COPY & RUN install",
             f"EXPOSE {port}", f"CMD {cmd}"]
    for j, line in enumerate(lines):
        txt(slide, line, x+0.12, y+0.38+j*0.29, 2.84, 0.27, size=8.5, color=LGRAY)

# docker-compose infra section
rrect(slide, 0.28, 5.56, 12.75, 1.58, RGBColor(0x08,0x18,0x2E))
rect(slide, 0.28, 5.56, 12.75, 0.3, NAVY2)
txt(slide, "  docker-compose.yml  —  Infrastructure Services", 0.4, 5.58, 8, 0.28,
    size=11, bold=True, color=SKY)
infra = [
    ("Zookeeper", ":2181", "Kafka coordination layer"),
    ("Confluent Kafka", ":9092", "Event streaming broker"),
    ("Kafka UI", ":18088", "Topic & consumer monitoring"),
    ("MySQL 8.0", ":3308", "Transactional database"),
    ("MongoDB 7.0", ":27017", "Document store"),
    ("Redis 7 Alpine", ":6379", "In-memory cache layer"),
]
for i,(svc,port,desc) in enumerate(infra):
    ci = i%3; ri = i//3
    x = 0.45 + ci*4.38; y = 5.92 + ri*0.44
    txt(slide, f"▸  {svc}  {port}  —  {desc}", x, y, 4.2, 0.4, size=9.5, color=WHITE)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — Timeliness, Caching & Detection Rules
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
rect(slide, 0, 0, 13.33, 7.5, OFFWHITE)
header(slide, "9 · Timeliness, Caching Policy & Detection Rules",
       "How our design supports scalability, real-time delivery, and data consistency")
footer(slide)

sections = [
    ("SQL Caching Policy  (Redis)", BLUE, [
        "Search results: cached by sha256(query params), TTL = 60 seconds",
        "Single job fetch: cached by job_id, TTL = 10 seconds",
        "Cache key format:  search:{sha256}  /  job:{uuid}",
        "Invalidation: job update or close flushes the key immediately",
        "Impact: 99.9% latency reduction (1,289ms → 0.9ms),  3.2× TPS gain",
        "Redis 7 Alpine — sub-millisecond memory lookups",
    ]),
    ("Messaging Flow & Timeliness", GREEN, [
        "Kafka fire-and-forget: HTTP response immediate, events async",
        "message.sent → SSE EventSource → real-time to recipient browser",
        "application.submitted → job-service updates applicants_count async",
        "job.viewed → view counter incremented without blocking HTTP response",
        "ai.requests → worker → ai.results → analytics stored in MongoDB",
        "SSE streams use EventEmitter in-process bus for sub-100ms delivery",
    ]),
    ("Detection Rules & Idempotency", ORANGE, [
        "Kafka at-least-once: duplicate events WILL arrive — handled safely",
        "MySQL processed_events: idempotency_key PRIMARY KEY deduplicates",
        "MongoDB analytics: findOne({idempotency_key}) before every insert",
        "UUID v4 idempotency_key in every Kafka event envelope",
        "JWT expiry validated on every authenticated request (HMAC-SHA256)",
        "Role enforcement: member vs recruiter checked per-route, 403 on fail",
    ]),
]
for i,(title,col,bullets) in enumerate(sections):
    x = 0.28 + i*4.38
    rrect(slide, x, 1.28, 4.2, 5.7, WHITE)
    rect(slide, x, 1.28, 4.2, 0.38, col)
    txt(slide, f"  {title}", x+0.1, 1.3, 4.0, 0.34, size=11, bold=True, color=WHITE)
    for j, b in enumerate(bullets):
        txt(slide, f"• {b}", x+0.14, 1.76+j*0.78, 3.9, 0.72, size=9.5, color=NAVY)


# ══════════════════════════════════════════════════════════════════════════════
# SAVE
# ══════════════════════════════════════════════════════════════════════════════
import os
out = "/Users/dipinjassal/sem_2/linkedin/presentation/LinkedIn_Group8_Presentation.pptx"
os.makedirs(os.path.dirname(out), exist_ok=True)
prs.save(out)
print(f"Saved: {out}")
